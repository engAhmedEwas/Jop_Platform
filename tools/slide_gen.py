#!/usr/bin/env python3
"""Generate a simple PowerPoint from docs/SLIDES.md.

Rules:
- Each top-level 'Slide X' section (lines starting with 'Slide') becomes a slide.
- The first non-empty line before 'Slide' is used as title if present.
"""
import sys
from pathlib import Path
import re


def parse_slides(md_text):
    # Split by lines that start with 'Slide' (case-insensitive)
    lines = md_text.splitlines()
    slides = []
    current = {"title": None, "bullets": []}
    for line in lines:
        m = re.match(r"^\s*Slide\s+\d+\s*[—-]\s*(.*)$", line, re.I)
        if m:
            # start a new slide
            if current["title"] or current["bullets"]:
                slides.append(current)
            current = {"title": m.group(1).strip() or None, "bullets": []}
            continue
        # ignore header lines
        if line.strip().startswith('#'):
            continue
        # bullet lines
        bm = re.match(r"^\s*[-*+]\s+(.*)$", line)
        if bm:
            current["bullets"].append(bm.group(1).strip())
            continue
        # short non-empty lines become subtitle if no title
        if line.strip() and not current["title"] and not line.strip().startswith('I can'):
            current["title"] = line.strip()

    if current["title"] or current["bullets"]:
        slides.append(current)
    return slides


def make_pptx(slides, outpath):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    for s in slides:
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        # title
        if slide.shapes.title and s.get('title'):
            slide.shapes.title.text = s.get('title')
        # body
        body = None
        for shp in slide.shapes:
            if shp.has_text_frame and shp is not slide.shapes.title:
                body = shp
                break
        if body is None:
            # add textbox
            left = Inches(1)
            top = Inches(1.8)
            width = Inches(8)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            body = txBox
        tf = body.text_frame
        tf.clear()
        for i, b in enumerate(s.get('bullets', [])):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            p.level = 0

    prs.save(outpath)


def main():
    repo = Path(__file__).resolve().parents[1]
    mdfile = repo / 'docs' / 'SLIDES.md'
    out = repo / 'docs' / 'slides.pptx'
    md = mdfile.read_text(encoding='utf-8')
    slides = parse_slides(md)
    if not slides:
        print('No slides parsed')
        sys.exit(2)
    make_pptx(slides, str(out))
    print('Saved', out)


if __name__ == '__main__':
    main()
